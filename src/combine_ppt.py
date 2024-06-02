import os
from pptx import Presentation

def extract_text_from_pptx(directory, output_file):
    all_text = []

    files = [f for f in os.listdir(directory) if f.endswith('.pptx')]

    for file in files:
        file_path = os.path.join(directory, file)
        presentation = Presentation(file_path)

        for slide in presentation.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            all_text.append(run.text)

    with open(output_file, 'w', encoding='utf-8') as f:
        f.write('\n'.join(all_text))

    print(f"Combined text saved as '{output_file}'")

output_file_path = os.path.join('d:\\work', 'combined_text.txt')
extract_text_from_pptx('d:\\work', output_file_path)
s